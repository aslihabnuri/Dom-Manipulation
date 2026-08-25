"""Render the real .pptx to HTML by reading every shape back out of the file."""
import base64, os, sys, html as H
from pptx import Presentation
from pptx.util import Emu
EMU=914400.0; DPI=96
def px(v): return (v or 0)/EMU*DPI
def col(c, dflt=None):
    try:
        if c is None or c.type is None: return dflt
        return "#%s" % str(c.rgb)
    except Exception: return dflt

prs=Presentation(sys.argv[1]); out=[]
SW, SH = px(prs.slide_width), px(prs.slide_height)
for si, sl in enumerate(prs.slides, 1):
    bg="#FFFFFF"
    try:
        f=sl.background.fill
        if f.type is not None and f.type==1: bg=col(f.fore_color,"#FFFFFF")
    except Exception: pass
    parts=[f'<div class="slide" style="width:{SW}px;height:{SH}px;background:{bg}"><div class="num">{si}</div>']
    for sh in sl.shapes:
        x,y,w,h = px(sh.left),px(sh.top),px(sh.width),px(sh.height)
        base=f"left:{x}px;top:{y}px;width:{w}px;height:{h}px;"
        if sh.shape_type==13 or sh.__class__.__name__=="Picture":
            img=sh.image; b=base64.b64encode(img.blob).decode()
            parts.append(f'<img style="position:absolute;{base}" src="data:{img.content_type};base64,{b}">'); continue
        if sh.has_table:
            t=sh.table
            rows=[]
            colw=[px(c.width) for c in t.columns]
            for ri,r in enumerate(t.rows):
                rowpx=px(r.height)
                cells=[]
                for ci,c in enumerate(r.cells):
                    cf=col(c.fill.fore_color,"#FFFFFF") if c.fill.type==1 else "transparent"
                    txt=""; st=""
                    for p in c.text_frame.paragraphs:
                        for run in p.runs:
                            fs=run.font.size.pt if run.font.size else 11
                            fc=col(run.font.color,"#282828")
                            fn=run.font.name or "Arimo"
                            bold="700" if run.font.bold else "400"
                            st=f"font-family:'{fn}';font-size:{fs}px;color:{fc};font-weight:{bold};"
                            txt+=H.escape(run.text)
                        al=p.alignment
                    cells.append(f'<td style="width:{colw[ci]}px;background:{cf};{st}text-align:'
                                 f'{"right" if txt.startswith("Rp") or txt=="0" else "left"}">{txt}</td>')
                rows.append(f'<tr style="height:{rowpx}px">'+"".join(cells)+"</tr>")
            rh=[px(r.height) for r in t.rows]
            realh=sum(rh)
            tb=f"left:{x}px;top:{y}px;width:{w}px;"
            parts.append(f'<table style="position:absolute;{tb}border-collapse:collapse">'+"".join(rows)+"</table>")
            continue
        fill="transparent"
        try:
            if sh.fill.type==1: fill=col(sh.fill.fore_color,"transparent")
        except Exception: pass
        rad="6px" if "ROUNDED" in str(sh.shape_type) else "0"
        parts.append(f'<div style="position:absolute;{base}background:{fill};border-radius:{rad}"></div>')
        if not sh.has_text_frame: continue
        tf=sh.text_frame
        ml=px(tf.margin_left); mr=px(tf.margin_right)
        va=str(tf.vertical_anchor or "")
        just="center" if "CENTER" in va or "MIDDLE" in va else "flex-start"
        ps=[]
        for p in tf.paragraphs:
            runs=[]
            for run in p.runs:
                fs=run.font.size.pt if run.font.size else 12
                fc=col(run.font.color,"#282828"); fn=run.font.name or "Arimo"
                bold="700" if run.font.bold else "400"
                sp=""
                try:
                    v=run.font._rPr.get('spc')
                    if v: sp=f"letter-spacing:{int(v)/100.0}px;"
                except Exception: pass
                runs.append(f'<span style="font-family:\'{fn}\',Arial;font-size:{fs}px;color:{fc};'
                            f'font-weight:{bold};{sp}">{H.escape(run.text)}</span>')
            if runs: ps.append("<div>"+"".join(runs)+"</div>")
        al=""
        try:
            if tf.paragraphs and str(tf.paragraphs[0].alignment or "").find("RIGHT")>=0: al="text-align:right;"
        except Exception: pass
        parts.append(f'<div style="position:absolute;{base}padding:0 {mr}px 0 {ml}px;display:flex;'
                     f'flex-direction:column;justify-content:{just};{al}overflow:visible;line-height:1.18">'
                     +"".join(ps)+"</div>")
    parts.append("</div>")
    out.append("".join(parts))
open("preview.html","w").write(
 "<style>body{margin:0;background:#8a8a8a;font-family:Arimo}"
 ".slide{position:relative;margin:14px auto;box-shadow:0 2px 12px rgba(0,0,0,.4);overflow:hidden}"
 ".num{position:absolute;right:-0px;top:-0px;background:#000;color:#fff;font:11px Arial;padding:2px 6px;z-index:99}"
 "td{padding:0 12px;border:0.5px solid #CCCCCC}</style>" + "".join(out))
print("preview.html written:", len(prs.slides), "slides")
