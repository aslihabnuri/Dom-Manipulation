"""Tables render at rowH x rows, not at the stored placeholder height.
Wrap every cell with the real TTF metrics and check it fits its row."""
from pptx import Presentation
from PIL import ImageFont
import sys

FONTS="/tmp/claude-0/-home-user-Dom-Manipulation/37162ed6-62c4-5584-966a-d15d826db4ae/scratchpad/fonts"
EMU=914400.0
def face(n,b):
    if n and "Zalando" in n: return f"{FONTS}/ZalandoSansExpanded-{'700' if b else '400'}.ttf"
    return f"{FONTS}/Arimo-{'700' if b else '400'}.ttf"
cache={}
def fnt(p,px):
    if (p,px) not in cache: cache[(p,px)]=ImageFont.truetype(p,px)
    return cache[(p,px)]
def wpt(t,n,sz,b,csp):
    if not t: return 0.0
    px=max(4,int(round(sz*4))); return fnt(face(n,b),px).getlength(t)/4.0 + csp*len(t)

def wrap_lines(text,n,sz,b,csp,avail_pt):
    words=text.split(); lines=1; cur=""
    for w in words:
        t=(cur+" "+w).strip()
        if wpt(t,n,sz,b,csp) <= avail_pt or not cur: cur=t
        else: lines+=1; cur=w
    return lines, wpt(cur,n,sz,b,csp)

prs=Presentation(sys.argv[1]); bad=[]
for si,slide in enumerate(prs.slides,1):
    others=[]
    for sh in slide.shapes:
        if not sh.has_table and sh.left is not None:
            others.append((sh.left/EMU,(sh.width or 0)/EMU,sh.top/EMU,(sh.height or 0)/EMU,
                           (sh.text_frame.text[:26] if sh.has_text_frame else "shape")))
    for sh in slide.shapes:
        if not sh.has_table: continue
        tb=sh.table
        colW=[c.width/EMU for c in tb.columns]
        rowH=[r.height/EMU for r in tb.rows]
        top=sh.top/EMU; left=sh.left/EMU
        rendered=sum(rowH)
        print(f"S{si} table  x={left:.2f} y={top:.2f}  stored h={(sh.height or 0)/EMU:.2f}in  "
              f"RENDERED h={rendered:.2f}in ({len(rowH)} rows)  bottom={top+rendered:.2f}in")
        if top+rendered > 7.5-0.5+0.02:
            bad.append(f"S{si}: table bottom {top+rendered:.2f}in past safe area (7.00in)")
        for ri,row in enumerate(tb.rows):
            for ci,cell in enumerate(row.cells):
                pad=((cell.margin_left or 0)+(cell.margin_right or 0))/EMU
                avail=(colW[ci]-pad)*72
                for p in cell.text_frame.paragraphs:
                    txt="".join(r.text for r in p.runs)
                    if not txt.strip(): continue
                    r0=p.runs[0]
                    sz=r0.font.size.pt if r0.font.size else 12
                    b=bool(r0.font.bold); n=r0.font.name
                    csp=0.0
                    ln,_=wrap_lines(txt,n,sz,b,csp,avail)
                    need=(ln*sz*1.22)/72.0 + 0.08
                    flag = need > rowH[ri]+0.005
                    if flag or ln>2:
                        bad.append(f"S{si} r{ri}c{ci}: '{txt[:44]}' -> {ln} line(s), "
                                   f"needs {need:.2f}in, row {rowH[ri]:.2f}in")
                    used=wpt(txt,n,sz,b,csp)
                    print(f"    r{ri}c{ci} {ln}L  w={used/72:.2f}/{avail/72:.2f}in "
                          f"({100*used/avail:.0f}%)  {'FIT' if not flag else 'OVER'}  {txt[:48]}")
        right=left+sum(colW)
        for ox,ow,oy,oh,lbl in others:
            vx=min(right,ox+ow)-max(left,ox)
            vy=min(top+rendered,oy+oh)-max(top,oy)
            if vx>0.03 and vy>0.03:
                bad.append(f"S{si}: '{lbl}' overlaps rendered table by {vx:.2f}x{vy:.2f}in")
print()
print(("TABLE ISSUES:\n  "+"\n  ".join(bad)) if bad else "tables: clean")
