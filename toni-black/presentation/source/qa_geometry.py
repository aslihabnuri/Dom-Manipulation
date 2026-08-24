"""QA the real .pptx: text fit, bounds and overlap, using the actual brand TTFs."""
from pptx import Presentation
from pptx.util import Emu
from PIL import ImageFont
import os, sys

FONTS="/tmp/claude-0/-home-user-Dom-Manipulation/37162ed6-62c4-5584-966a-d15d826db4ae/scratchpad/fonts"
def face(name, bold):
    if name and "Zalando" in name:
        return f"{FONTS}/ZalandoSansExpanded-{'700' if bold else '400'}.ttf"
    return f"{FONTS}/Arimo-{'700' if bold else '400'}.ttf"

EMU_IN=914400.0
SLIDE_W, SLIDE_H = 13.3, 7.5
MARGIN = 0.5
cache={}
def fnt(path, px):
    k=(path,px)
    if k not in cache: cache[k]=ImageFont.truetype(path, px)
    return cache[k]

def run_width_pt(text, fontname, size_pt, bold, charsp_pt):
    if not text: return 0.0
    px=max(4,int(round(size_pt*4)))          # 4x for precision, scale back after
    f=fnt(face(fontname,bold), px)
    w=f.getlength(text)/4.0
    return w + charsp_pt*len(text)

prs=Presentation(sys.argv[1])
issues=[]
for si,slide in enumerate(prs.slides,1):
    boxes=[]
    for sh in slide.shapes:
        if sh.left is None: continue
        x,y=sh.left/EMU_IN, sh.top/EMU_IN
        w,h=(sh.width or 0)/EMU_IN,(sh.height or 0)/EMU_IN
        full_bleed = (w>=SLIDE_W-0.01 and h>=SLIDE_H-0.01)
        if not full_bleed:
            if x < MARGIN-0.02 or y < MARGIN-0.02 or x+w > SLIDE_W-MARGIN+0.02 or y+h > SLIDE_H-MARGIN+0.02:
                issues.append(f"S{si}: shape outside safe margin  x={x:.2f} y={y:.2f} w={w:.2f} h={h:.2f}")
        if not sh.has_text_frame: continue
        tf=sh.text_frame
        txt="".join(r.text for p in tf.paragraphs for r in p.runs)
        if not txt.strip(): continue
        boxes.append((x,y,w,h,txt[:28]))
        # read the real internal inset rather than assuming one
        try:
            pad = ((tf.margin_left or 0) + (tf.margin_right or 0)) / EMU_IN / 2.0
        except Exception:
            pad = 0.10
        for p in tf.paragraphs:
            line=""; sz=12; bold=False; fname=None; csp=0.0
            wsum=0.0
            for r in p.runs:
                sz = r.font.size.pt if r.font.size else 12
                bold = bool(r.font.bold)
                fname = r.font.name
                wsum += run_width_pt(r.text, fname, sz, bold, 0.0)
                line += r.text
            if not line.strip(): continue
            avail_pt=(w-2*pad)*72
            # single-line elements (short labels, numbers) must not exceed the box
            if "\n" not in line and wsum > avail_pt*1.02 and len(line) < 60:
                issues.append(f"S{si}: text may overflow  '{line[:40]}'  needs {wsum/72:.2f}in, box {w:.2f}in")
    # overlap between text boxes on the same slide
    for i in range(len(boxes)):
        for j in range(i+1,len(boxes)):
            a,b=boxes[i],boxes[j]
            ox=min(a[0]+a[2],b[0]+b[2])-max(a[0],b[0])
            oy=min(a[1]+a[3],b[1]+b[3])-max(a[1],b[1])
            if ox>0.05 and oy>0.05:
                issues.append(f"S{si}: text boxes overlap  '{a[4]}' x '{b[4]}'  ({ox:.2f}x{oy:.2f}in)")
print(f"slides: {len(prs.slides)}")
if issues:
    print(f"ISSUES ({len(issues)}):")
    for i in issues: print("  -",i)
else:
    print("no geometry issues found")
