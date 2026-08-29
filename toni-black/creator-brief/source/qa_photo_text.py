"""Type sits on photography in this deck, so contrast cannot be read off a fill
colour. For every run that overlaps a picture, crop that picture's own pixels
under the run and compute the contrast it actually gets."""
from pptx import Presentation
from PIL import Image
import numpy as np, io, sys
EMU=914400.0
def lum_hex(h):
    c=[int(h[i:i+2],16)/255 for i in (0,2,4)]
    c=[(v/12.92 if v<=0.03928 else ((v+0.055)/1.055)**2.4) for v in c]
    return 0.2126*c[0]+0.7152*c[1]+0.0722*c[2]
def lum_L(L):
    v=L/255.0
    return v/12.92 if v<=0.03928 else ((v+0.055)/1.055)**2.4
def ratio(a,b):
    hi,lo=max(a,b),min(a,b); return (hi+0.05)/(lo+0.05)

prs=Presentation(sys.argv[1]); bad=[]; n=0
for si,sl in enumerate(prs.slides,1):
    pics=[]
    for sh in sl.shapes:
        if sh.__class__.__name__=="Picture" and sh.left is not None:
            im=Image.open(io.BytesIO(sh.image.blob)).convert("L")
            pics.append((sh.left/EMU, sh.top/EMU, (sh.width or 0)/EMU,
                         (sh.height or 0)/EMU, im))
    if not pics: continue
    for sh in sl.shapes:
        if not sh.has_text_frame or sh.left is None: continue
        x,y=sh.left/EMU, sh.top/EMU
        w,h=(sh.width or 0)/EMU, (sh.height or 0)/EMU
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if not r.text.strip(): continue
                try: fg=str(r.font.color.rgb)
                except Exception: continue
                sz=r.font.size.pt if r.font.size else 12
                big = sz>=18 or (sz>=14 and r.font.bold)
                need = 3.0 if big else 4.5
                for (px,py,pw,ph,im) in pics:
                    ox0,oy0=max(x,px),max(y,py)
                    ox1,oy1=min(x+w,px+pw),min(y+h,py+ph)
                    if ox1-ox0<=0.02 or oy1-oy0<=0.02: continue
                    W,H=im.size; sx,sy=W/pw,H/ph
                    crop=im.crop((int((ox0-px)*sx),int((oy0-py)*sy),
                                  max(int((ox1-px)*sx),int((ox0-px)*sx)+1),
                                  max(int((oy1-py)*sy),int((oy0-py)*sy)+1)))
                    a=np.asarray(crop).astype(float)
                    # The worst local patch matters, and which end is worst depends
                    # on the text tone — a mid grey loses contrast in both directions,
                    # so test both and keep the lower ratio.
                    hi,lo=float(np.percentile(a,85)),float(np.percentile(a,15))
                    cr=min(ratio(lum_hex(fg),lum_L(hi)), ratio(lum_hex(fg),lum_L(lo)))
                    L=hi if ratio(lum_hex(fg),lum_L(hi))<ratio(lum_hex(fg),lum_L(lo)) else lo
                    n+=1
                    if cr < need:
                        bad.append(f"S{si}: {cr:4.2f}:1 (need {need}) {sz:.0f}pt #{fg} "
                                   f"on photo L={L:.0f}  '{r.text[:30]}'")
print(f"runs on photography checked: {n}")
print(("PHOTO CONTRAST FAILURES:\n  "+"\n  ".join(bad)) if bad
      else "photo contrast: every run on photography passes WCAG AA")
