"""QA the real .pptx: text fit, bounds and overlap, using the actual brand TTFs."""
from pptx import Presentation
from pptx.util import Emu
from PIL import ImageFont
import os, sys

import os
FONTS=os.path.normpath(os.path.join(os.path.dirname(os.path.abspath(__file__)),
    "..", "..", "september-nine-to-nine", "source", "fonts"))
def face(name, bold):
    if name and "Zalando" in name:
        return f"{FONTS}/ZalandoSansExpanded-{'700' if bold else '400'}.ttf"
    return f"{FONTS}/Arimo-{'700' if bold else '400'}.ttf"

EMU_IN=914400.0
SLIDE_W = SLIDE_H = None   # read from the file, not assumed
MARGIN = 0.30
cache={}
def fnt(path, px):
    k=(path,px)
    if k not in cache: cache[k]=ImageFont.truetype(path, px)
    return cache[k]

def run_width_pt(text, fontname, size_pt, bold, charsp_pt):
    if not text: return 0.0
    px=max(4,int(round(size_pt*4)))          # 4x for precision, scale back after
    f=fnt(face(fontname,bold), px)
    w=f.getlength(text)/4.0
    return w + charsp_pt*len(text)

prs=Presentation(sys.argv[1])
SLIDE_W, SLIDE_H = prs.slide_width/EMU_IN, prs.slide_height/EMU_IN
issues=[]; notes=[]
for si,slide in enumerate(prs.slides,1):
    boxes=[]
    pics=[]
    for sh in slide.shapes:                       # collect pictures first
        if sh.__class__.__name__=="Picture" and sh.left is not None:
            pics.append((sh.left/EMU_IN, sh.top/EMU_IN,
                         (sh.width or 0)/EMU_IN, (sh.height or 0)/EMU_IN))
    for sh in slide.shapes:
        if sh.left is None: continue
        x,y=sh.left/EMU_IN, sh.top/EMU_IN
        w,h=(sh.width or 0)/EMU_IN,(sh.height or 0)/EMU_IN
        is_pic = sh.__class__.__name__=="Picture"
        full_bleed = (w>=SLIDE_W-0.01 and h>=SLIDE_H-0.01)
        if not full_bleed and not is_pic:      # photographs are meant to bleed
            if x < MARGIN-0.02 or y < MARGIN-0.02 or x+w > SLIDE_W-MARGIN+0.02 or y+h > SLIDE_H-MARGIN+0.02:
                issues.append(f"S{si}: shape outside safe margin  x={x:.2f} y={y:.2f} w={w:.2f} h={h:.2f}")
        if not sh.has_text_frame: continue
        tf=sh.text_frame
        txt="".join(r.text for p in tf.paragraphs for r in p.runs)
        if not txt.strip(): continue
        boxes.append((x,y,w,h,txt[:28]))
        # read the real internal inset rather than assuming one
        try:
            pad = ((tf.margin_left or 0) + (tf.margin_right or 0)) / EMU_IN / 2.0
        except Exception:
            pad = 0.10
        avail_pt=(w-2*pad)*72
        used_in=0.0
        for p in tf.paragraphs:
            line=""; sz=12; bold=False; fname=None; csp=0.0
            for r in p.runs:
                sz = r.font.size.pt if r.font.size else 12
                bold = bool(r.font.bold); fname = r.font.name
                rPr = r._r.rPr
                if rPr is not None and rPr.get("spc"): csp = int(rPr.get("spc"))/100.0
                line += r.text
            if not line.strip():
                used_in += sz*1.2/72.0; continue
            # a single word wider than the box cannot wrap — that is a hard overflow
            for word in line.split():
                if run_width_pt(word, fname, sz, bold, csp) > avail_pt*1.02:
                    issues.append(f"S{si}: word too wide  '{word}' in '{line[:32]}'  "
                                  f"needs {run_width_pt(word,fname,sz,bold,csp)/72:.2f}in, "
                                  f"box {w-2*pad:.2f}in")
                    break
            # wrap it the way a renderer would, then measure the height it needs
            lines, cur = 1, ""
            for word in line.split():
                t=(cur+" "+word).strip()
                if run_width_pt(t, fname, sz, bold, csp) <= avail_pt or not cur: cur=t
                else: lines+=1; cur=word
            ls = p.line_spacing
            lh = (ls.pt if hasattr(ls,"pt") else (float(ls)*sz if ls else sz*1.2))
            used_in += lines*lh/72.0
        if used_in > h + 0.02:
            issues.append(f"S{si}: text taller than its box  '{txt[:34]}'  "
                          f"needs {used_in:.2f}in, box {h:.2f}in")
    # text sitting on top of a picture — nothing in this deck is designed that way
    for (x,y,w,h,lbl) in boxes:
        for (px,py,pw,ph) in pics:
            ox=min(x+w,px+pw)-max(x,px); oy=min(y+h,py+ph)-max(y,py)
            if ox>0.05 and oy>0.05:
                notes.append(f"S{si}: text on photography  '{lbl}'")
    # overlap between text boxes on the same slide
    for i in range(len(boxes)):
        for j in range(i+1,len(boxes)):
            a,b=boxes[i],boxes[j]
            ox=min(a[0]+a[2],b[0]+b[2])-max(a[0],b[0])
            oy=min(a[1]+a[3],b[1]+b[3])-max(a[1],b[1])
            if ox>0.05 and oy>0.05:
                issues.append(f"S{si}: text boxes overlap  '{a[4]}' x '{b[4]}'  ({ox:.2f}x{oy:.2f}in)")
print(f"slides: {len(prs.slides)}")
print(f"text placed on photography: {len(notes)} runs (contrast checked by qa_photo_text.py)")
if issues:
    print(f"ISSUES ({len(issues)}):")
    for i in issues: print("  -",i)
else:
    print("no geometry issues found")
