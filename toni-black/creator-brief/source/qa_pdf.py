"""Every string in the .pptx must survive into the .pdf, on the same page."""
import re, sys, pymupdf
from pptx import Presentation

def norm(t): return re.sub(r"\s+", " ", t).strip()

prs = Presentation(sys.argv[1]); doc = pymupdf.open(sys.argv[2])
if len(prs.slides) != doc.page_count:
    print(f"FAIL page count {doc.page_count} != slides {len(prs.slides)}"); sys.exit(1)
missing = []
for i, sl in enumerate(prs.slides):
    want = []
    for sh in sl.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            want += [norm(x) for x in sh.text_frame.text.split("\n") if norm(x)]
        if sh.has_table:
            want += [norm(c.text) for r in sh.table.rows for c in r.cells if norm(c.text)]
    got = norm(doc[i].get_text().replace("\n", " "))
    got_ns = got.replace(" ", "")
    for wtxt in want:
        # letter-spaced runs come back with spaces between glyphs; compare both ways
        if wtxt not in got and wtxt.replace(" ", "") not in got_ns:
            missing.append(f"p{i+1}: {wtxt!r}")
    print(f"page {i+1}: {len(want)} strings checked")
b = doc[0].rect
print(f"page size {b.width/72:.3f} x {b.height/72:.3f} in")
print(("MISSING:\n  " + "\n  ".join(missing)) if missing else "all pptx text present in pdf")
