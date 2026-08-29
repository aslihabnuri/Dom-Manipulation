"""WCAG contrast for every run, measured against the shape actually behind it."""
from pptx import Presentation
import sys
EMU=914400.0
def lum(h):
    c=[int(h[i:i+2],16)/255 for i in (0,2,4)]
    c=[(v/12.92 if v<=0.03928 else ((v+0.055)/1.055)**2.4) for v in c]
    return 0.2126*c[0]+0.7152*c[1]+0.0722*c[2]
def ratio(a,b):
    la,lb=lum(a),lum(b); hi,lo=max(la,lb),min(la,lb)
    return (hi+0.05)/(lo+0.05)
def fill_of(sh):
    try:
        if sh.fill.type==1 and sh.fill.fore_color.type is not None:
            return str(sh.fill.fore_color.rgb)
    except Exception: pass
    return None

prs=Presentation(sys.argv[1]); bad=[]; n=0
for si,sl in enumerate(prs.slides,1):
    bg="FFFFFF"
    try:
        f=sl.background.fill
        if f.type==1 and f.fore_color.type is not None: bg=str(f.fore_color.rgb)
    except Exception: pass
    grounds=[]                                   # filled shapes, in z-order
    for sh in sl.shapes:
        if sh.left is None: continue
        c=fill_of(sh)
        if c: grounds.append((sh.left/EMU,sh.top/EMU,(sh.width or 0)/EMU,(sh.height or 0)/EMU,c))
    for sh in sl.shapes:
        if not sh.has_text_frame or sh.left is None: continue
        x,y=sh.left/EMU,sh.top/EMU
        cx,cy=x+(sh.width or 0)/EMU/2, y+(sh.height or 0)/EMU/2
        ground=bg
        for (gx,gy,gw,gh,gc) in grounds:         # last one wins = topmost
            if gx-0.01<=cx<=gx+gw+0.01 and gy-0.01<=cy<=gy+gh+0.01: ground=gc
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if not r.text.strip(): continue
                try: fg=str(r.font.color.rgb)
                except Exception: continue
                sz=r.font.size.pt if r.font.size else 12
                big = sz>=18 or (sz>=14 and r.font.bold)
                need = 3.0 if big else 4.5
                cr=ratio(fg,ground); n+=1
                if cr < need:
                    bad.append(f"S{si}: {cr:4.2f}:1 (need {need}) {sz:.0f}pt "
                               f"#{fg} on #{ground}  '{r.text[:34]}'")
print(f"runs checked: {n}")
print(("CONTRAST FAILURES:\n  "+"\n  ".join(bad)) if bad else "contrast: all runs pass WCAG AA")
