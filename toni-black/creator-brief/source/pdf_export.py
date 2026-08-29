"""Render the real .pptx to PDF: read every shape back out of the file and
lay it out at true scale (positions in inches, type in points), then print.

LibreOffice cannot load any document in this environment — its import filters
are broken — so this is the conversion path.
"""
import base64, os, sys, math, asyncio, html as H
from pptx import Presentation
from playwright.async_api import async_playwright

SP  = os.path.dirname(os.path.abspath(__file__))
FD  = os.path.normpath(os.path.join(os.path.dirname(os.path.abspath(__file__)),
    "..", "..", "september-nine-to-nine", "source", "fonts"))
EMU = 914400.0
A   = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
C   = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"

def IN(v): return (v or 0) / EMU

def col(c, dflt=None):
    try:
        if c is None or c.type is None: return dflt
        return "#%s" % str(c.rgb)
    except Exception:
        return dflt

def font_css():
    out = []
    for fam, weight, f in [("Zalando Sans Expanded", 400, "ZalandoSansExpanded-600"),
                           ("Zalando Sans Expanded", 700, "ZalandoSansExpanded-700"),
                           ("Arimo", 400, "Arimo-400"), ("Arimo", 700, "Arimo-700")]:
        b = base64.b64encode(open(f"{FD}/{f}.ttf", "rb").read()).decode()
        out.append(f"@font-face{{font-family:'{fam}';font-weight:{weight};font-style:normal;"
                   f"src:url(data:font/ttf;base64,{b}) format('truetype');}}")
    return "".join(out)

ANCHOR = {"TOP": "flex-start", "MIDDLE": "center", "BOTTOM": "flex-end"}
ALIGN  = {"LEFT": "left", "CENTER": "center", "RIGHT": "right", "JUSTIFY": "justify"}

def runs_html(p, dflt_color="#282828"):
    out = []
    for r in p.runs:
        fs   = r.font.size.pt if r.font.size else 12
        fc   = col(r.font.color, dflt_color) or dflt_color
        fn   = r.font.name or "Arimo"
        bold = 700 if r.font.bold else 400
        spc  = ""
        rPr  = r._r.rPr
        if rPr is not None and rPr.get("spc"):
            spc = f"letter-spacing:{int(rPr.get('spc'))/100.0}pt;"
        # PowerPoint renders runs of consecutive spaces; HTML would collapse them
        out.append(f'<span style="font-family:\'{fn}\',sans-serif;font-size:{fs}pt;'
                   f'color:{fc};font-weight:{bold};{spc}white-space:pre-wrap">'
                   f'{H.escape(r.text)}</span>')
    return "".join(out)

def textframe_html(tf, w, h, dflt_color="#282828"):
    va = ANCHOR.get(str(tf.vertical_anchor).split(" ")[0], "flex-start")
    ml, mr = IN(tf.margin_left), IN(tf.margin_right)
    mt, mb = IN(tf.margin_top), IN(tf.margin_bottom)
    ps = []
    for p in tf.paragraphs:
        if not p.runs:
            ps.append("<div>&nbsp;</div>"); continue
        al = ALIGN.get(str(p.alignment).split(" ")[0], "left")
        # line_spacing is a Length (EMU-valued int subclass) for spcPts, a float
        # multiplier for spcPct. Printing the raw Length gives EMU, not points.
        lsv = p.line_spacing
        if lsv is None:            ls = "line-height:1.18;"
        elif hasattr(lsv, "pt"):   ls = f"line-height:{lsv.pt}pt;"
        else:                      ls = f"line-height:{float(lsv):.3f};"
        ps.append(f'<div style="text-align:{al};{ls}">{runs_html(p, dflt_color)}</div>')
    return (f'<div style="position:absolute;inset:0;display:flex;flex-direction:column;'
            f'justify-content:{va};padding:{mt}in {mr}in {mb}in {ml}in;">' + "".join(ps) + "</div>")

def donut_svg(chart, w_in, h_in):
    """Draw the doughnut from the chart part's own values and colours."""
    x = chart._chartSpace
    vals = [float(v.text) for v in x.iter(f"{C}v")
            if v.getparent().getparent().tag == f"{C}numCache"]
    cols = []
    for dpt in x.iter(f"{C}dPt"):
        s = dpt.find(f".//{A}srgbClr")
        if s is not None: cols.append("#" + s.get("val"))
    hole = x.find(f".//{C}holeSize")
    hole = int(hole.get("val")) / 100.0 if hole is not None else 0.5
    if not vals: return ""
    total = sum(vals)
    R, r = 50.0, 50.0 * hole
    cx = cy = 50.0
    segs, a0 = [], -math.pi / 2                      # 12 o'clock, clockwise
    for i, v in enumerate(vals):
        a1 = a0 + 2 * math.pi * v / total
        large = 1 if (a1 - a0) > math.pi else 0
        p = (f"M {cx+R*math.cos(a0):.4f} {cy+R*math.sin(a0):.4f} "
             f"A {R} {R} 0 {large} 1 {cx+R*math.cos(a1):.4f} {cy+R*math.sin(a1):.4f} "
             f"L {cx+r*math.cos(a1):.4f} {cy+r*math.sin(a1):.4f} "
             f"A {r} {r} 0 {large} 0 {cx+r*math.cos(a0):.4f} {cy+r*math.sin(a0):.4f} Z")
        segs.append(f'<path d="{p}" fill="{cols[i] if i < len(cols) else "#818284"}"/>')
        a0 = a1
    return ('<svg viewBox="0 0 100 100" style="position:absolute;inset:0;width:100%;height:100%">'
            + "".join(segs) + "</svg>")

def build(path):
    prs = Presentation(path)
    SW, SH = IN(prs.slide_width), IN(prs.slide_height)
    pages = []
    for sl in prs.slides:
        bg = "#FFFFFF"
        try:
            f = sl.background.fill
            if f.type is not None and f.type == 1: bg = col(f.fore_color, "#FFFFFF")
        except Exception: pass
        parts = [f'<div class="pg" style="width:{SW}in;height:{SH}in;background:{bg}">']
        for sh in sl.shapes:
            x, y, w, h = IN(sh.left), IN(sh.top), IN(sh.width), IN(sh.height)
            box = f"position:absolute;left:{x}in;top:{y}in;width:{w}in;height:{h}in;"

            if sh.__class__.__name__ == "Picture":
                img = sh.image
                b = base64.b64encode(img.blob).decode()
                parts.append(f'<img style="{box}" src="data:{img.content_type};base64,{b}">')
                continue

            if sh.has_table:
                t = sh.table
                colw = [IN(c.width) for c in t.columns]
                rows = []
                for r in t.rows:
                    cells = []
                    for ci, c in enumerate(r.cells):
                        try:
                            cf = col(c.fill.fore_color, "transparent") if c.fill.type == 1 else "transparent"
                        except Exception:
                            cf = "transparent"
                        ml, mr = IN(c.margin_left), IN(c.margin_right)
                        inner = []
                        for p in c.text_frame.paragraphs:
                            if not p.runs: continue
                            al = ALIGN.get(str(p.alignment).split(" ")[0], "left")
                            inner.append(f'<div style="text-align:{al};line-height:1.18">{runs_html(p)}</div>')
                        cells.append(f'<td style="width:{colw[ci]}in;background:{cf};'
                                     f'padding:0 {mr}in 0 {ml}in;vertical-align:middle;'
                                     f'border:0.5pt solid #CCCCCC">' + "".join(inner) + "</td>")
                    rows.append(f'<tr style="height:{IN(r.height)}in">' + "".join(cells) + "</tr>")
                parts.append(f'<table style="position:absolute;left:{x}in;top:{y}in;width:{w}in;'
                             f'border-collapse:collapse;table-layout:fixed">' + "".join(rows) + "</table>")
                continue

            if sh.__class__.__name__ == "GraphicFrame" and getattr(sh, "has_chart", False):
                parts.append(f'<div style="{box}">' + donut_svg(sh.chart, w, h) + "</div>")
                continue

            fill = "transparent"
            try:
                if sh.fill.type == 1: fill = col(sh.fill.fore_color, "transparent")
            except Exception: pass
            line = ""
            try:
                lc = col(sh.line.color, None)
                lw = sh.line.width.pt if sh.line.width else 0
                if lc and lw: line = f"border:{lw}pt solid {lc};"
            except Exception: pass
            rad = "0"
            try:
                adj = list(sh.adjustments)
                if adj: rad = f"{adj[0]*min(w,h):.4f}in"
            except Exception: pass
            inner = ""
            if sh.has_text_frame and sh.text_frame.text.strip():
                inner = textframe_html(sh.text_frame, w, h)
            parts.append(f'<div style="{box}background:{fill};{line}border-radius:{rad};'
                         f'box-sizing:border-box">{inner}</div>')
        parts.append("</div>")
        pages.append("".join(parts))

    return (f'<meta charset="utf-8"><style>{font_css()}'
            f'@page{{size:{SW}in {SH}in;margin:0}}'
            f'*{{margin:0;padding:0;box-sizing:border-box}}'
            f'html,body{{-webkit-print-color-adjust:exact;print-color-adjust:exact}}'
            f'.pg{{position:relative;overflow:hidden;page-break-after:always}}'
            f'.pg:last-child{{page-break-after:auto}}'
            f'</style>' + "".join(pages), SW, SH)

async def main():
    src = sys.argv[1]; dst = sys.argv[2]
    html, SW, SH = build(src)
    hp = os.path.join(SP, "_pdf.html")
    open(hp, "w", encoding="utf-8").write(html)
    async with async_playwright() as pw:
        b = await pw.chromium.launch(
            executable_path="/opt/pw-browsers/chromium-1194/chrome-linux/chrome")
        pg = await b.new_page()
        await pg.goto("file://" + hp)
        await pg.wait_for_timeout(1500)
        await pg.pdf(path=dst, width=f"{SW}in", height=f"{SH}in",
                     print_background=True, margin={"top":"0","bottom":"0","left":"0","right":"0"})
        await b.close()
    print("pdf:", dst, os.path.getsize(dst), "bytes")

asyncio.run(main())
